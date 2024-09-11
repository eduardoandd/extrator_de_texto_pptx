from pptx import Presentation

arquivo='BDA03_Functions_Procedures.pptx'

presentation=Presentation(arquivo)

slide=presentation.slides[17]

for shape in slide.shapes:
    if hasattr(shape,'text'):
        print(shape.text)